from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()

    # Define colors for theme
    bg_color = RGBColor(245, 245, 245) # Light gray background
    title_color = RGBColor(44, 62, 80) # Dark Blue
    text_color = RGBColor(51, 51, 51) # Dark Gray
    accent_color = RGBColor(52, 152, 219) # Light Blue
    sidebar_color = RGBColor(41, 128, 185) # Strong Blue

    # Helper function to apply background design
    def apply_background_design(slide):
        # 1. Set solid background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # 2. Add a decorative sidebar
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.4)
        height = prs.slide_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = sidebar_color
        shape.line.fill.background() # No border

    # Helper function to set font for a text frame
    def set_font(text_frame, font_size=Pt(20), is_bold=False, color=text_color):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.size = font_size
                run.font.color.rgb = color
                run.font.bold = is_bold
    
    # Same as set_font but applied while creating content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        apply_background_design(slide)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.name = 'Times New Roman'
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True

        # Content
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear existing placeholder content

        for text in content_text_list:
            p = tf.add_paragraph()
            p.text = text
            p.font.name = 'Times New Roman'
            p.font.size = Pt(24)
            p.font.color.rgb = text_color
            p.space_after = Pt(14)
            p.level = 0
            
    # 1. Title Slide Customization
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    apply_background_design(slide)
    
    title = slide.shapes.title
    title.text = "Certificate Management System"
    title.text_frame.paragraphs[0].font.name = 'Times New Roman'
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "A Web Application for College Tech Events"
    subtitle.text_frame.paragraphs[0].font.name = 'Times New Roman'
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color

    # 2. About the Project
    add_slide("1. About the Project", [
        "The Certificate Management System is a web-based application designed to automate the process of issuing certificates.",
        "It serves as a centralized platform for college tech events.",
        "Administrators can easily generate and issue certificates to participants.",
        "Participants can log in to view, download, and request corrections for their certificates.",
        "The goal is to reduce paperwork and make certificate distribution efficient and eco-friendly."
    ])

    # 3. Hardware & Software Specifications
    add_slide("2. Hardware & Software Specifications", [
        "Hardware Requirements:",
        "    - RAM: 2 GB",
        "    - ROM: 256 GB",
        "    - Processor: Compatible with Windows 7 (32-bit)",
        "",
        "Software Requirements:",
        "    - OS: Windows 7 (32-bit)",
        "    - Server Environment: XAMPP v3.2.2",
        "    - Language: PHP v7.3.2",
        "    - Database: MySQL v5.10"
    ])

    # 4. Problem Definition
    add_slide("3. Problem Definition", [
        "Constraint of Time: Manual writing or printing of certificates takes a lot of time.",
        "Distribution Issues: Distributing physical certificates to hundreds of students is difficult.",
        "Human Errors: Spelling mistakes in names involve re-printing, leading to waste.",
        "Record Maintenance: Keeping track of who received a certificate on paper is hard to manage.",
        "Lack of Access: Students cannot access their certificates digitally if they lose the hard copy."
    ])

    # 5. System Study
    add_slide("4. System Study", [
        "Existing System:",
        "- Entirely manual process.",
        "- Organizers collect names manually and print certificates.",
        "- High cost of printing and paper.",
        "- No backup for lost certificates.",
        "",
        "Proposed System:",
        "- Fully automated and digital.",
        "- Database stores all records permanently.",
        "- Instant generation and distribution.",
        "- Cost-effective and paperless."
    ])

    # 6. Proposed System
    add_slide("5. Proposed System", [
        "The proposed system is a Web Application built with PHP and MySQL.",
        "It provides role-based access for Admins and Users.",
        "Admins can bulk-generate certificates for events.",
        "Users can securely login to download their certificates in PDF format.",
        "Includes a correction request system for fixing name errors before final printing.",
        "Ensures data integrity and easy accessibility."
    ])

    # 7. Dataflow Diagrams (Digitally Drawn)
    
    # Helpers for drawing DFDs
    def draw_entity(slide, left, top, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(1.5), Inches(0.8))
        shape.text = text
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_color
        shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].font.bold = True
        return shape

    def draw_process(slide, left, top, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(1.5), Inches(1.5))
        shape.text = text
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
        shape.text_frame.paragraphs[0].font.size = Pt(12)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        return shape

    def draw_store(slide, left, top, text):
        # Using a slightly different shape for Data Store, e.g., Flowchart Manual Operation or similar, 
        # but often just an open rectangle. We'll use a shape that looks distinct.
        shape = slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, Inches(1.5), Inches(0.8))
        shape.text = text
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
        shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
        shape.text_frame.paragraphs[0].font.size = Pt(12)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        return shape

    def draw_arrow(slide, shape_from, shape_to, text=""):
        # Simple connector
        connector = slide.shapes.add_connector(1, shape_from.left + shape_from.width/2, shape_from.top + shape_from.height/2,
                                              shape_to.left + shape_to.width/2, shape_to.top + shape_to.height/2)
        connector.line.width = Pt(2)
        if text:
            # Add a text label near the midpoint manually (simplified)
            mid_x = (shape_from.left + shape_to.left) / 2
            mid_y = (shape_from.top + shape_to.top) / 2
            tb = slide.shapes.add_textbox(mid_x, mid_y, Inches(1), Inches(0.4))
            tb.text = text
            tb.text_frame.paragraphs[0].font.size = Pt(10)
            tb.text_frame.paragraphs[0].font.name = 'Times New Roman'

    # Level 0
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    apply_background_design(slide)
    slide.shapes.title.text = "6. DFD Level 0"
    slide.shapes.title.text_frame.paragraphs[0].font.name = 'Times New Roman'
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    admin = draw_entity(slide, Inches(1), Inches(3), "Admin")
    system = draw_process(slide, Inches(4.25), Inches(2.75), "Certificate\nSystem")
    student = draw_entity(slide, Inches(7.5), Inches(3), "Student")

    draw_arrow(slide, admin, system, "Manage Certs")
    draw_arrow(slide, student, system, "Login/Download")

    # Level 1
    slide = prs.slides.add_slide(slide_layout)
    apply_background_design(slide)
    slide.shapes.title.text = "DFD Level 1"
    slide.shapes.title.text_frame.paragraphs[0].font.name = 'Times New Roman'
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

    # Processes
    login_proc = draw_process(slide, Inches(1), Inches(2), "1.0 Login")
    manage_proc = draw_process(slide, Inches(4), Inches(2), "2.0 Manage\nCertificates")
    download_proc = draw_process(slide, Inches(4), Inches(4.5), "3.0 Download\nCertificate")
    correct_proc = draw_process(slide, Inches(7), Inches(3.5), "4.0 Request\nCorrection")

    # DBs
    user_db = draw_store(slide, Inches(1), Inches(4.5), "Users DB")
    cert_db = draw_store(slide, Inches(7), Inches(2), "Certs DB")

    # Connections
    draw_arrow(slide, login_proc, user_db)
    draw_arrow(slide, login_proc, manage_proc)
    draw_arrow(slide, manage_proc, cert_db)
    draw_arrow(slide, cert_db, download_proc)
    draw_arrow(slide, download_proc, correct_proc)
    
    # Level 2 (Correction Flow)
    slide = prs.slides.add_slide(slide_layout)
    apply_background_design(slide)
    slide.shapes.title.text = "DFD Level 2 (Correction)"
    slide.shapes.title.text_frame.paragraphs[0].font.name = 'Times New Roman'
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

    user_ent = draw_entity(slide, Inches(0.5), Inches(3), "User")
    sub_proc = draw_process(slide, Inches(2.5), Inches(3), "4.1 Submit\nRequest")
    val_proc = draw_process(slide, Inches(4.5), Inches(3), "4.2 Validate")
    update_proc = draw_process(slide, Inches(6.5), Inches(3), "4.3 Update\nStatus")
    req_db = draw_store(slide, Inches(4.5), Inches(5), "Requests DB")
    
    draw_arrow(slide, user_ent, sub_proc)
    draw_arrow(slide, sub_proc, val_proc)
    draw_arrow(slide, val_proc, update_proc)
    draw_arrow(slide, val_proc, req_db)

    prs.save('Certificate_Management_System_v4.pptx')
    print("Presentation created successfully.")

if __name__ == "__main__":
    create_presentation()
